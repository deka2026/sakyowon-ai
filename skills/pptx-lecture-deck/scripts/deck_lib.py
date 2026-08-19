# -*- coding: utf-8 -*-
"""강의안 PPTX 공통 라이브러리 (python-pptx)

사용법:
    from deck_lib import *
    prs = new_deck()
    slide_question(prs, 1, "질문?", "한 줄 답", [[("제목", ["불릿", "*강조불릿"])], [...]],
                   footer="자료: ...")
    prs.save(out)

핵심은 card_height()/fit_font() — 한글 줄 수를 추정해 카드 높이를 잡고,
열이 넘치면 본문 폰트를 단계적으로 줄인다. 렌더 검증은 export_png.ps1로.
"""
import math

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------- 팔레트 (필요하면 프로젝트별로 교체) ----------
DARK = RGBColor(0x14, 0x4A, 0x38)
MID = RGBColor(0x2E, 0x8B, 0x6B)
LIGHT = RGBColor(0xE8, 0xF2, 0xEE)
LIME = RGBColor(0x8D, 0xC6, 0x3F)
INK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x6E, 0x7B, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARDBG = RGBColor(0xF7, 0xF9, 0xF8)
LINE = RGBColor(0xD5, 0xE2, 0xDC)
WARN = RGBColor(0xC0, 0x5A, 0x2B)

FONT = "맑은 고딕"          # run마다 지정해야 한다. 테마에 맡기면 영문 폰트로 렌더됨
SW, SH = 13.333, 7.5        # 16:9

# 본문 레이아웃 상수
MARGIN = 0.55
BODY_W = 12.23
BODY_TOP = 2.15
BODY_BOTTOM = 6.85          # 각주가 있을 때
CARD_GAP = 0.20
COL_GAP = 0.26


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------- 기본 도형/텍스트 ----------
def para(tf, text, size=12, bold=False, color=INK, space_after=4,
         align=PP_ALIGN.LEFT, first=False, bullet=None, line=1.15, space_before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    r = p.add_run()
    r.text = f"{bullet} {text}" if bullet else text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    s.shadow.inherit = False        # 기본 그림자를 끄지 않으면 카드가 지저분해진다
    s.text_frame.word_wrap = True
    return s


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


# ---------- 넘침 방지 계산 ----------
def est_lines(text, size, width_in):
    """한글 기준 줄 수 추정. 0.88은 실측 보정값(0.80은 과소추정 → 넘침)."""
    char_w = size * 0.88 / 72.0
    per_line = max(8, int(width_in / char_w))
    return max(1, math.ceil(len(text) / per_line))


def card_height(heading, bullets, size, cw):
    """카드 하나에 필요한 높이(inch). 제목 글자도 본문과 함께 줄어든다."""
    tw = cw - 0.42
    hs = min(13.0, size + 2.5)
    h = 0.16
    h += est_lines(heading, hs, cw - 0.4) * (hs * 1.20 / 72.0) + 0.15
    for b in bullets:
        t = b[1:] if b.startswith("*") else b
        h += est_lines("▸ " + t, size, tw) * (size * 1.22 / 72.0) + 4.5 / 72.0
    h += 0.22
    return h


def fit_font(cols, cw, avail_of, candidates=(11.0, 10.5, 10.0, 9.5, 9.0, 8.5)):
    """열별 필요 높이가 가용 높이에 들어갈 때까지 폰트를 낮춘다.
    끝까지 안 들어가면 최소값을 돌려주지만, 그때는 폰트가 아니라 문장을 줄일 것."""
    size = candidates[-1]
    for cand in candidates:
        if all(sum(card_height(h, b, cand, cw) for h, b in col) <= avail_of(ci, len(col))
               for ci, col in enumerate(cols)):
            return cand
        size = cand
    return size


# ---------- 카드형 질문 슬라이드 ----------
def slide_question(prs, num, question, answer, cols, footer=None,
                   warn=False, image_col=None):
    """cols: [[(제목, [불릿...]), ...], [...]]   불릿 앞 '*'는 강조
    image_col: (열 인덱스, 이미지경로, 높이inch) — 그 열의 카드 아래에 이미지 배치"""
    s = blank(prs)

    rect(s, 0, 0, SW, 1.12, fill=DARK)
    rect(s, MARGIN, 0.24, 0.66, 0.64, fill=LIME)
    tb = textbox(s, MARGIN, 0.28, 0.66, 0.56)
    para(tb.text_frame, f"Q{num}", size=17, bold=True, color=DARK,
         first=True, align=PP_ALIGN.CENTER)
    tb = textbox(s, 1.42, 0.22, 11.4, 0.75)
    para(tb.text_frame, question, size=23.5, bold=True, color=WHITE, first=True, line=1.05)

    band = RGBColor(0xFB, 0xF0, 0xE8) if warn else LIGHT
    rect(s, 0, 1.12, SW, 0.78, fill=band)
    rect(s, 0, 1.12, 0.12, 0.78, fill=WARN if warn else MID)
    tb = textbox(s, MARGIN, 1.20, 12.3, 0.65)
    para(tb.text_frame, answer, size=14.5, bold=True,
         color=WARN if warn else DARK, first=True, line=1.15)

    top, bottom = BODY_TOP, (BODY_BOTTOM if footer else 7.08)
    n = len(cols)
    cw = (BODY_W - COL_GAP * (n - 1)) / n

    def avail_of(ci, ncards):
        a = (bottom - top) - CARD_GAP * (ncards - 1)
        if image_col and image_col[0] == ci:
            a -= image_col[2] + CARD_GAP
        return a

    size = fit_font(cols, cw, avail_of)

    x = MARGIN
    for ci, col in enumerate(cols):
        avail = avail_of(ci, len(col))
        needs = [card_height(h, b, size, cw) for h, b in col]
        tot = sum(needs)
        if tot > avail:
            heights = [nd * avail / tot for nd in needs]      # 최후의 수단
        else:
            extra = avail - tot
            heights = [nd + extra * (nd / tot) for nd in needs]
        y = top
        for (heading, bullets), h in zip(col, heights):
            rect(s, x, y, cw, h, fill=CARDBG, line_color=LINE)
            rect(s, x, y, cw, 0.05, fill=MID)
            hs = min(13.0, size + 2.5)
            tb = textbox(s, x + 0.2, y + 0.11, cw - 0.4, 0.4)
            para(tb.text_frame, heading, size=hs, bold=True, color=DARK, first=True)
            tb = textbox(s, x + 0.2, y + 0.11 + hs * 1.20 / 72.0 + 0.13,
                         cw - 0.38, h - 0.4)
            tf = tb.text_frame
            for i, b in enumerate(bullets):
                strong = b.startswith("*")
                para(tf, b[1:] if strong else b, size=size,
                     color=DARK if strong else INK, bold=strong,
                     first=(i == 0), space_after=4.5, bullet="▸", line=1.22)
            y += h + CARD_GAP
        if image_col and image_col[0] == ci:
            _, path, ih = image_col
            from PIL import Image
            pw, ph = Image.open(path).size
            iw = ih * pw / ph
            s.shapes.add_picture(path, Inches(x + (cw - iw) / 2), Inches(y + 0.02),
                                 height=Inches(ih))
        x += cw + COL_GAP

    if footer:
        rect(s, MARGIN, 6.92, BODY_W, 0.03, fill=LINE)
        tb = textbox(s, MARGIN, 6.98, BODY_W, 0.4)
        para(tb.text_frame, footer, size=9.5, color=GRAY, first=True)
    return s
